from __future__ import annotations

from pathlib import Path
from typing import Protocol, cast

import pptx as python_pptx
from pptx.shapes.base import BaseShape
from pptx.slide import Slide
from pptx.text.text import TextFrame

from src.core.constants import CHUNK_SECTION_KEY
from src.core.exceptions import DocumentLoadError
from src.domain.entities.document import Document


class _TextFrameShape(Protocol):
    has_text_frame: bool
    text_frame: TextFrame


def shape_text(shape: BaseShape) -> str:
    if not shape.has_text_frame:
        return ""
    text_frame = cast(_TextFrameShape, shape).text_frame
    paragraphs = [p.text for p in text_frame.paragraphs]
    return "\n".join(p for p in paragraphs if p.strip())


def slide_title(slide: Slide) -> str | None:
    title_shape = slide.shapes.title
    if title_shape is None:
        return None
    title = shape_text(title_shape).strip()
    return title or None


def slide_text(slide: Slide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        text = shape_text(shape)
        if text.strip():
            parts.append(text)
    return "\n\n".join(parts)


class PptxLoader:
    """Loads text from PPTX files using python-pptx."""

    @staticmethod
    def load(path: Path) -> Document:
        try:
            presentation = python_pptx.Presentation(str(path))

            # Per-slide records keep title↔body alignment for SectionChunker.
            # "sections" remain named titles only (omits untitled slides).
            slide_records: list[dict[str, object]] = []
            section_titles: list[str] = []
            for slide in presentation.slides:
                title = slide_title(slide)
                text = slide_text(slide)
                if not text.strip():
                    continue
                slide_records.append({"title": title, "text": text})
                if title:
                    section_titles.append(title)

            content = "\n\n---\n\n".join(
                str(record["text"]) for record in slide_records
            )

            metadata: dict[str, object] = {
                "filename": path.name,
                "extension": path.suffix.lower(),
                "loader": "pptx",
                "slide_count": len(presentation.slides),
                "sections": section_titles,
                "slides": slide_records,
            }
            if section_titles:
                metadata[CHUNK_SECTION_KEY] = section_titles[0]

            return Document(
                source=str(path.resolve()),
                content=content,
                metadata=metadata,
            )
        except DocumentLoadError:
            raise
        except Exception as exc:
            raise DocumentLoadError(f"Cannot load PPTX: {path}", cause=exc) from exc
