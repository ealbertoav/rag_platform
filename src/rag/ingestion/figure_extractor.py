"""Figure asset extraction and Chunk builders (T-230).

When "parsing.figure_assets.enabled" is true, persist figure bytes from
layout "figures[]" (Docling PDF/DOCX) or PPTX picture shapes under a
:class:`~src.rag.ingestion.local_asset_store.LocalAssetStore`, then set
"figures[].asset_path" and produce :class:`~src.domain.entities.chunk.Chunk`
instances with "asset_path" / "figure_id".

Docling picture export reuses the layout parser's picture-enabled
:class:`~docling.document_converter.DocumentConverter` so PDF/DOCX
"figures[]" and rasterized "pictures" share the same pipeline options.
Pictures are matched to layout entries by page and bbox provenance when
present, then by remaining document order — not by blind list index —
so a second conversion cannot silently attach the wrong asset.

DOCX falls back to embedded "python-docx" image parts when Docling
conversion fails, Docling is unavailable, or "PictureItem.get_image()"
returns no bytes. Fallback blobs are collected from body "a:blip"
embeds in document reading order (not "part.rels" iteration order),
excluding header/footer relationships. PPTX walks picture shapes
recursively through group shapes.

Soft-fails per figure when bytes cannot be exported. VLM captions are
T-231; caption indexing is T-232.
"""

from __future__ import annotations

import io
import logging
import uuid
from pathlib import Path
from typing import Any, Protocol, cast

from src.core.constants import (
    ASSET_PATH_KEY,
    BBOX_KEY,
    CHUNK_PAGE_KEY,
    CHUNK_SOURCE_KEY,
    CHUNK_TYPE_FIGURE,
    CHUNK_TYPE_KEY,
    FIGURE_ID_KEY,
    MODALITY_FIGURE,
)
from src.core.exceptions import ConfigurationError, DocumentLoadError
from src.core.settings import FigureAssetSettings, Settings
from src.domain.entities.chunk import Chunk
from src.domain.entities.document import Document
from src.rag.chunking.metadata import chunk_metadata
from src.rag.ingestion.local_asset_store import LocalAssetStore, document_asset_key

logger = logging.getLogger(__name__)

_PPTX_EXTENSION = ".pptx"
_DOCLING_EXTENSIONS: frozenset[str] = frozenset({".pdf", ".docx"})
_FIGURE_CHUNK_NAMESPACE = uuid.UUID("c4e91b7a-2d6f-4a18-9e3b-8f0d5c1a7b42")
_DEFAULT_FIGURE_TEXT = "[figure]"


class _PilImage(Protocol):
    def save(self, fp: Any, *args: Any, **kwargs: Any) -> None: ...


class _PptxImage(Protocol):
    blob: bytes
    ext: str | None


class _PictureShape(Protocol):
    image: _PptxImage


class _DoclingPicture(Protocol):
    prov: list[Any]

    def caption_text(self, doc: object) -> str: ...

    def get_image(self, doc: object) -> _PilImage | None: ...


class _DoclingDocument(Protocol):
    pictures: list[_DoclingPicture]


class _DoclingConversionStatus(Protocol):
    name: str


class _DoclingConversionResult(Protocol):
    status: _DoclingConversionStatus | None
    document: _DoclingDocument


class _DoclingConverter(Protocol):
    def convert(self, source: str) -> _DoclingConversionResult: ...


def figure_chunk_id(source: str, figure_id: str) -> str:
    """Stable chunk ID for a figure asset at a *source*."""
    return str(uuid.uuid5(_FIGURE_CHUNK_NAMESPACE, f"{source}:{figure_id}"))


def is_figure_chunk(chunk: Chunk) -> bool:
    """Return True when *chunk* is a figure asset index point."""
    return chunk.metadata.get(CHUNK_TYPE_KEY) == CHUNK_TYPE_FIGURE


def build_figure_chunks(document: Document) -> list[Chunk]:
    """Build figure chunks from document metadata that already has asset paths."""
    figures = document.metadata.get("figures")
    if not isinstance(figures, list) or not figures:
        return []

    chunks: list[Chunk] = []
    for index, raw_entry in enumerate(figures):
        if not isinstance(raw_entry, dict):
            logger.debug("Skipping non-dict figure entry at index %d", index)
            continue
        figure_id = raw_entry.get(FIGURE_ID_KEY)
        asset_path = raw_entry.get(ASSET_PATH_KEY) or raw_entry.get("asset_path")
        if not figure_id or not asset_path:
            continue

        caption = raw_entry.get("caption")
        if isinstance(caption, str) and caption.strip():
            text = caption.strip()
        else:
            text = _DEFAULT_FIGURE_TEXT

        metadata = chunk_metadata(document.metadata)
        metadata[CHUNK_TYPE_KEY] = CHUNK_TYPE_FIGURE
        metadata[FIGURE_ID_KEY] = str(figure_id)
        metadata[ASSET_PATH_KEY] = str(asset_path)
        metadata[CHUNK_SOURCE_KEY] = document.source
        if CHUNK_PAGE_KEY in raw_entry:
            metadata[CHUNK_PAGE_KEY] = raw_entry[CHUNK_PAGE_KEY]
        if BBOX_KEY in raw_entry:
            metadata[BBOX_KEY] = raw_entry[BBOX_KEY]

        chunks.append(
            Chunk(
                id=figure_chunk_id(document.source, str(figure_id)),
                document_id=document.id,
                text=text,
                metadata=metadata,
                modality=MODALITY_FIGURE,
                asset_path=str(asset_path),
            )
        )
    return chunks


def apply_figure_assets(
    document: Document,
    path: Path,
    *,
    app_settings: Settings | None = None,
    store: LocalAssetStore | None = None,
) -> Document:
    """Persist figure bytes and attach "asset_path" onto "figures[]" metadata.

    No-op when figure assets are disabled. Soft-fails (logs and returns the
    original document) on configuration / conversion errors, so ingest continues.
    """
    cfg = _figure_asset_settings(app_settings)
    if not cfg.enabled:
        return document

    asset_store = store or LocalAssetStore(cfg.store_dir)
    try:
        suffix = path.suffix.lower()
        if suffix == _PPTX_EXTENSION:
            return _apply_pptx_figures(document, path, asset_store)
        if suffix in _DOCLING_EXTENSIONS:
            return _apply_docling_figures(document, path, asset_store)
        logger.debug("Skipping figure assets for unsupported type %s", path.name)
        return document
    except ConfigurationError as exc:
        logger.warning(
            "Figure asset extraction misconfigured for %s: %s — continuing without assets",
            path.name,
            exc,
        )
        return document
    except DocumentLoadError as exc:
        logger.warning(
            "Figure asset extraction failed for %s: %s — continuing without assets",
            path.name,
            exc,
        )
        return document
    except Exception as exc:
        logger.warning(
            "Unexpected figure asset error for %s: %s — continuing without assets",
            path.name,
            exc,
        )
        return document


def _figure_asset_settings(app_settings: Settings | None) -> FigureAssetSettings:
    if app_settings is None:
        from src.core.settings import settings as default_settings

        return default_settings.parsing.figure_assets
    return app_settings.parsing.figure_assets


def _figure_entry_needs_asset(entry: dict[str, Any]) -> bool:
    """Return True when a figures[] dict has an id but no persisted asset_path."""
    if not entry.get(FIGURE_ID_KEY):
        return False
    asset_path = entry.get(ASSET_PATH_KEY) or entry.get("asset_path")
    return not asset_path


def _apply_pptx_figures(
    document: Document,
    path: Path,
    store: LocalAssetStore,
) -> Document:
    extracted = _extract_pptx_picture_bytes(path)
    existing = document.metadata.get("figures")
    if not extracted:
        pending = 0
        if isinstance(existing, list):
            pending = sum(
                1
                for entry in existing
                if isinstance(entry, dict) and _figure_entry_needs_asset(entry)
            )
        if pending:
            logger.warning(
                "PPTX %s has %d figure(s) without asset_path but no extractable pictures — "
                "continuing without assets",
                path.name,
                pending,
            )
        return document

    doc_key = document_asset_key(document.source)
    figures: list[dict[str, Any]] = []
    existing_by_id: dict[str, dict[str, Any]] = {}
    if isinstance(existing, list):
        for entry in existing:
            if isinstance(entry, dict) and entry.get(FIGURE_ID_KEY):
                existing_by_id[str(entry[FIGURE_ID_KEY])] = dict(entry)

    for index, (blob, extension, page) in enumerate(extracted, start=1):
        figure_id = f"figure-{index}"
        entry = dict(existing_by_id.get(figure_id, {FIGURE_ID_KEY: figure_id}))
        entry[FIGURE_ID_KEY] = figure_id
        entry[CHUNK_PAGE_KEY] = page
        try:
            asset_path = store.save(doc_key, figure_id, blob, extension=extension)
            entry[ASSET_PATH_KEY] = str(asset_path)
        except Exception as exc:
            logger.warning(
                "Failed to store PPTX figure %s from %s: %s",
                figure_id,
                path.name,
                exc,
            )
        figures.append(entry)

    metadata = dict(document.metadata)
    metadata["figures"] = figures
    return document.model_copy(update={"metadata": metadata})


def _iter_pptx_shapes(shapes: Any) -> Any:
    """Yield shapes depth-first, descending into group shapes."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        shape_type = getattr(shape, "shape_type", None)
        if shape_type == MSO_SHAPE_TYPE.GROUP:
            nested = getattr(shape, "shapes", None)
            if nested is not None:
                yield from _iter_pptx_shapes(nested)
            continue
        yield shape


def _extract_pptx_picture_bytes(path: Path) -> list[tuple[bytes, str, int]]:
    try:
        import pptx as python_pptx
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as exc:
        raise ConfigurationError(
            "PPTX figure extraction requires python-pptx. Install with: uv pip install python-pptx"
        ) from exc

    try:
        presentation = python_pptx.Presentation(str(path))
    except Exception as exc:
        raise DocumentLoadError(
            f"Cannot open PPTX for figure extraction: {path}",
            cause=exc,
        ) from exc

    pictures: list[tuple[bytes, str, int]] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in _iter_pptx_shapes(slide.shapes):
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                image = cast(_PictureShape, shape).image
                blob = image.blob
                extension = (image.ext or "png").lstrip(".").lower() or "png"
            except Exception as exc:
                logger.warning(
                    "Skipping unreadable PPTX picture on slide %s in %s: %s",
                    slide_index,
                    path.name,
                    exc,
                )
                continue
            if blob:
                pictures.append((blob, extension, slide_index))
    return pictures


def _apply_docling_figures(
    document: Document,
    path: Path,
    store: LocalAssetStore,
) -> Document:
    existing = document.metadata.get("figures")
    if not isinstance(existing, list) or not existing:
        logger.debug("No layout figures[] for %s — skipping Docling asset export", path.name)
        return document

    figure_entries = [entry for entry in existing if isinstance(entry, dict)]
    expected_count = len(figure_entries)
    pictures = _load_docling_pictures_or_empty(path)
    docx_blobs = _docx_fallback_blobs(path)

    if not pictures and not docx_blobs:
        logger.warning(
            "Layout figures[] has %d entr%s for %s but Docling returned no pictures — "
            "continuing without asset_path",
            expected_count,
            "y" if expected_count == 1 else "ies",
            path.name,
        )
        return document

    # Docling pictures are primary (matched by page+bbox, then remaining order).
    # DOCX embedded body images backfill missing rasters by figure slot index.
    available = max(len(pictures), len(docx_blobs)) if docx_blobs else len(pictures)
    if expected_count > available:
        if docx_blobs and not pictures:
            logger.warning(
                "DOCX embedded images returned %d picture(s) for %s but layout figures[] "
                "has %d dict entr%s — unmatched figures leave without asset_path",
                len(docx_blobs),
                path.name,
                expected_count,
                "y" if expected_count == 1 else "ies",
            )
        else:
            logger.warning(
                "Docling returned %d picture(s) for %s but layout figures[] has %d dict entr%s — "
                "unmatched figures leave without asset_path",
                len(pictures),
                path.name,
                expected_count,
                "y" if expected_count == 1 else "ies",
            )

    doc_key = document_asset_key(document.source)
    figures: list[dict[str, Any]] = []
    used_picture_indexes: set[int] = set()
    slot = 0
    for raw_entry in existing:
        if not isinstance(raw_entry, dict):
            continue
        entry = dict(raw_entry)
        figure_id = str(entry.get(FIGURE_ID_KEY) or f"figure-{len(figures) + 1}")
        entry[FIGURE_ID_KEY] = figure_id

        blob: bytes | None = None
        extension = "png"
        attempted_docling = False
        picture_index = _match_picture_index(entry, pictures, used_picture_indexes, slot)
        if picture_index is not None:
            attempted_docling = True
            used_picture_indexes.add(picture_index)
            picture, doc = pictures[picture_index]
            try:
                blob = _picture_to_png_bytes(picture, doc)
            except Exception as exc:
                logger.warning(
                    "Failed to export Docling figure %s from %s: %s",
                    figure_id,
                    path.name,
                    exc,
                )

        if not blob and slot < len(docx_blobs):
            blob, extension = docx_blobs[slot]
            logger.debug(
                "Using python-docx embedded image for %s in %s",
                figure_id,
                path.name,
            )

        slot += 1

        if not blob:
            if attempted_docling:
                logger.warning(
                    "No image bytes for %s in %s — leaving without asset_path",
                    figure_id,
                    path.name,
                )
            else:
                logger.warning(
                    "No Docling picture for %s in %s — leaving without asset_path",
                    figure_id,
                    path.name,
                )
            figures.append(entry)
            continue

        try:
            asset_path = store.save(doc_key, figure_id, blob, extension=extension)
            entry[ASSET_PATH_KEY] = str(asset_path)
        except Exception as exc:
            logger.warning(
                "Failed to export Docling figure %s from %s: %s",
                figure_id,
                path.name,
                exc,
            )
        figures.append(entry)

    metadata = dict(document.metadata)
    metadata["figures"] = figures
    return document.model_copy(update={"metadata": metadata})


def _picture_page(picture: _DoclingPicture) -> int | None:
    prov = getattr(picture, "prov", None)
    if not prov:
        return None
    try:
        return int(prov[0].page_no)
    except (AttributeError, IndexError, TypeError, ValueError):
        return None


def _picture_bbox(picture: _DoclingPicture) -> list[float] | None:
    prov = getattr(picture, "prov", None)
    if not prov:
        return None
    try:
        box = prov[0].bbox
    except (AttributeError, IndexError, TypeError):
        return None
    if box is None:
        return None
    try:
        return [float(box.l), float(box.t), float(box.r), float(box.b)]
    except (AttributeError, TypeError, ValueError):
        return None


def _bboxes_equal(left: list[float], right: list[float], *, tol: float = 1e-3) -> bool:
    if len(left) != 4 or len(right) != 4:
        return False
    return all(abs(a - b) <= tol for a, b in zip(left, right, strict=True))


def _match_picture_index(
    entry: dict[str, Any],
    pictures: list[tuple[_DoclingPicture, _DoclingDocument]],
    used: set[int],
    fallback_index: int,
) -> int | None:
    """Map a layout figure entry to a Docling picture index.

    Prefer page + bbox provenance equality so a second Docling conversion with a
    different picture order still attaches the correct raster. Fall back to the
    next unused document-order slot.
    """
    if not pictures:
        return None

    page = entry.get(CHUNK_PAGE_KEY)
    bbox = entry.get(BBOX_KEY)
    if page is not None and isinstance(bbox, list) and len(bbox) == 4:
        try:
            page_int = int(page)
            bbox_floats = [float(value) for value in bbox]
        except (TypeError, ValueError):
            page_int = None
            bbox_floats = None
        if page_int is not None and bbox_floats is not None:
            for index, (picture, _) in enumerate(pictures):
                if index in used:
                    continue
                if _picture_page(picture) != page_int:
                    continue
                picture_bbox = _picture_bbox(picture)
                if picture_bbox is not None and _bboxes_equal(picture_bbox, bbox_floats):
                    return index

    if fallback_index < len(pictures) and fallback_index not in used:
        return fallback_index
    for index in range(len(pictures)):
        if index not in used:
            return index
    return None


def _docx_fallback_blobs(path: Path) -> list[tuple[bytes, str]]:
    """Return embedded DOCX image blobs for figure-asset fallback, else []."""
    if path.suffix.lower() != ".docx":
        return []
    try:
        return _extract_docx_picture_bytes(path)
    except ConfigurationError:
        raise
    except DocumentLoadError as exc:
        logger.warning(
            "DOCX embedded-image fallback unavailable for %s: %s",
            path.name,
            exc,
        )
        return []
    except Exception as exc:
        logger.warning(
            "DOCX embedded-image fallback failed for %s: %s",
            path.name,
            exc,
        )
        return []


def _extension_from_image_content_type(content_type: str) -> str:
    normalized = (content_type or "").lower().strip()
    mapping = {
        "image/png": "png",
        "image/jpeg": "jpg",
        "image/jpg": "jpg",
        "image/gif": "gif",
        "image/bmp": "bmp",
        "image/tiff": "tiff",
        "image/webp": "webp",
        "image/x-emf": "emf",
        "image/x-wmf": "wmf",
    }
    if normalized in mapping:
        return mapping[normalized]
    if "/" in normalized:
        subtype = normalized.rsplit("/", 1)[-1]
        if subtype.startswith("x-"):
            subtype = subtype[2:]
        if subtype:
            return subtype
    return "png"


def _extract_docx_picture_bytes(path: Path) -> list[tuple[bytes, str]]:
    """Extract body-embedded image parts from a DOCX in document reading order.

    Walks "a:blip" embeds under the document body, so order matches layout
    figures and header/footer relationships on "part.rels" are ignored.
    Repeated "rId" references are kept as separate slots (the same bytes may
    backfill multiple figure entries).
    """
    try:
        import docx as python_docx
        from docx.oxml.ns import qn
    except ImportError as exc:
        raise ConfigurationError(
            "DOCX figure fallback requires python-docx. Install with: uv pip install python-docx"
        ) from exc

    try:
        document = python_docx.Document(str(path))
    except Exception as exc:
        raise DocumentLoadError(
            f"Cannot open DOCX for figure extraction: {path}",
            cause=exc,
        ) from exc

    body = getattr(getattr(document, "element", None), "body", None)
    if body is None:
        return []

    embed_attr = qn("r:embed")
    blip_tag = qn("a:blip")
    rels = document.part.rels
    pictures: list[tuple[bytes, str]] = []
    for blip in body.iter(blip_tag):
        rid = blip.get(embed_attr)
        if not rid:
            continue
        try:
            rel = rels[rid]
        except KeyError:
            logger.warning(
                "Skipping DOCX image with missing relationship %s in %s",
                rid,
                path.name,
            )
            continue
        reltype = str(getattr(rel, "reltype", "") or "")
        if "image" not in reltype.lower():
            continue
        try:
            target = rel.target_part
            blob = bytes(getattr(target, "blob", b"") or b"")
            content_type = str(getattr(target, "content_type", "") or "")
            extension = _extension_from_image_content_type(content_type)
        except Exception as exc:
            logger.warning(
                "Skipping unreadable DOCX embedded image in %s: %s",
                path.name,
                exc,
            )
            continue
        if not blob:
            continue
        pictures.append((blob, extension))
    return pictures


def _load_docling_pictures_or_empty(
    path: Path,
) -> list[tuple[_DoclingPicture, _DoclingDocument]]:
    """Load Docling pictures, soft-failing so DOCX embedded-image fallback can run.

    "DocumentLoadError" always soft-fails to an empty list. "ConfigurationError"
    (e.g., missing Docling) soft-fails only for ".docx", where python-docx can still
    supply bytes; for other types it re-raises so ingest reports misconfiguration.
    """
    try:
        return _load_docling_pictures(path)
    except DocumentLoadError as exc:
        logger.warning(
            "Docling figure conversion failed for %s: %s — "
            "trying embedded-image fallback if available",
            path.name,
            exc,
        )
        return []
    except ConfigurationError:
        if path.suffix.lower() != ".docx":
            raise
        logger.warning(
            "Docling unavailable for %s — trying python-docx embedded-image fallback",
            path.name,
        )
        return []


def _load_docling_pictures(path: Path) -> list[tuple[_DoclingPicture, _DoclingDocument]]:
    converter = _create_picture_converter()
    try:
        result = converter.convert(str(path))
    except ConfigurationError:
        raise
    except Exception as exc:
        raise DocumentLoadError(
            f"Cannot convert {path.name} for figure asset extraction",
            cause=exc,
        ) from exc

    status = getattr(result, "status", None)
    if status is not None and getattr(status, "name", str(status)) == "FAILURE":
        raise DocumentLoadError(f"Docling conversion failed for figure assets: {path}")

    doc = result.document
    pictures = getattr(doc, "pictures", None) or []
    return [(picture, doc) for picture in pictures]


def _picture_to_png_bytes(picture: _DoclingPicture, doc: _DoclingDocument) -> bytes | None:
    image = picture.get_image(doc)
    if image is None:
        return None
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    data = buffer.getvalue()
    return data or None


def _create_picture_converter() -> _DoclingConverter:
    """Build the same picture-enabled Docling converter used by the layout parser."""
    from src.infrastructure.parsers.docling_parser import create_docling_converter

    return cast(_DoclingConverter, create_docling_converter())
