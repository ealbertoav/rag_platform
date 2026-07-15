"""T-010 — document loader tests.

PDF is tested via mocking (pypdf is a unit boundary, not under test here).
DOCX, HTML, and Markdown use real fixture files created in tmp_path.
"""

from __future__ import annotations

from pathlib import Path
from typing import Protocol, cast
from unittest.mock import MagicMock, patch

import docx as python_docx
import pptx as python_pptx
import pytest
from pptx.slide import Slide
from pptx.text.text import TextFrame
from pptx.util import Inches

from src.core.constants import CHUNK_SECTION_KEY
from src.core.exceptions import DocumentLoadError
from src.core.slide_records import SlideRecord
from src.domain.entities.document import Document
from src.domain.entities.parsed_document import ParsedDocument
from src.infrastructure.loaders import load_document
from src.infrastructure.loaders.docx_loader import DocxLoader
from src.infrastructure.loaders.html_loader import HtmlLoader
from src.infrastructure.loaders.markdown_loader import MarkdownLoader
from src.infrastructure.loaders.pdf_loader import PdfLoader
from src.infrastructure.loaders.pptx_loader import (
    PptxLoader,
    shape_text,
    slide_text,
    slide_title,
)

# ── Fixtures ───────────────────────────────────────────────────────────────────


class _TextFrameShape(Protocol):
    has_text_frame: bool
    text_frame: TextFrame


def _set_shape_text(shape: object | None, text: str) -> None:
    if shape is None:
        raise RuntimeError("shape is missing")
    typed = cast(_TextFrameShape, shape)
    assert typed.has_text_frame
    typed.text_frame.text = text


def _set_slide_title_and_body(slide: Slide, title: str, body: str) -> None:
    _set_shape_text(slide.shapes.title, title)
    _set_shape_text(slide.placeholders[1], body)


@pytest.fixture
def docx_file(tmp_path: Path) -> Path:
    path = tmp_path / "sample.docx"
    doc = python_docx.Document()
    doc.add_heading("Introduction", level=1)
    doc.add_paragraph("This is the first paragraph.")
    doc.add_heading("Details", level=2)
    doc.add_paragraph("This is the second paragraph with more detail.")
    doc.save(str(path))
    return path


@pytest.fixture
def pptx_file(tmp_path: Path) -> Path:
    path = tmp_path / "sample.pptx"
    prs = python_pptx.Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    _set_slide_title_and_body(title_slide, "Introduction", "Welcome to the deck.")

    content_slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_title_and_body(content_slide, "Details", "Second slide with more detail.")
    prs.save(str(path))
    return path


@pytest.fixture
def blank_pptx_file(tmp_path: Path) -> Path:
    path = tmp_path / "blank.pptx"
    prs = python_pptx.Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(str(path))
    return path


@pytest.fixture
def html_file(tmp_path: Path) -> Path:
    path = tmp_path / "sample.html"
    path.write_text(
        """<!DOCTYPE html>
<html>
<head><title>Test Page</title><style>body{color:red}</style></head>
<body>
  <nav>Navigation links</nav>
  <header>Site Header</header>
  <main>
    <h1>Main Content</h1>
    <p>This is the main paragraph.</p>
    <p>And a second one.</p>
  </main>
  <footer>Footer text</footer>
  <script>alert('hello')</script>
</body>
</html>""",
        encoding="utf-8",
    )
    return path


@pytest.fixture
def markdown_file(tmp_path: Path) -> Path:
    path = tmp_path / "sample.md"
    path.write_text(
        """# Title

First paragraph with some content.

## Section One

Content of section one.

## Section Two

Content of section two.
""",
        encoding="utf-8",
    )
    return path


@pytest.fixture
def latin1_file(tmp_path: Path) -> Path:
    path = tmp_path / "latin1.md"
    path.write_bytes("# Héllo Wörld\n\nSome cöntent.".encode("latin-1"))
    return path


# ── PdfLoader ──────────────────────────────────────────────────────────────────


class TestPdfLoader:
    @staticmethod
    def _make_mock_reader(page_texts: list[str]) -> MagicMock:
        pages = []
        for text in page_texts:
            page = MagicMock()
            page.extract_text.return_value = text
            pages.append(page)
        reader = MagicMock()
        reader.pages = pages
        return reader

    def test_returns_document(self, tmp_path: Path):
        path = tmp_path / "test.pdf"
        path.write_bytes(b"%PDF-1.4 placeholder")
        reader = self._make_mock_reader(["Page one content.", "Page two content."])
        with patch("src.infrastructure.loaders.pdf_loader.PdfReader", return_value=reader):
            doc = PdfLoader().load(path)
        assert isinstance(doc, Document)

    def test_content_joins_pages(self, tmp_path: Path):
        path = tmp_path / "test.pdf"
        path.write_bytes(b"%PDF-1.4 placeholder")
        reader = self._make_mock_reader(["First page.", "Second page."])
        with patch("src.infrastructure.loaders.pdf_loader.PdfReader", return_value=reader):
            doc = PdfLoader().load(path)
        assert "First page." in doc.content
        assert "Second page." in doc.content

    def test_metadata_has_required_keys(self, tmp_path: Path):
        path = tmp_path / "test.pdf"
        path.write_bytes(b"%PDF-1.4 placeholder")
        reader = self._make_mock_reader(["text"])
        with patch("src.infrastructure.loaders.pdf_loader.PdfReader", return_value=reader):
            doc = PdfLoader().load(path)
        assert doc.metadata["filename"] == "test.pdf"
        assert doc.metadata["extension"] == ".pdf"
        assert doc.metadata["loader"] == "pdf"
        assert doc.metadata["page_count"] == 1

    def test_source_is_absolute_path(self, tmp_path: Path):
        path = tmp_path / "test.pdf"
        path.write_bytes(b"%PDF-1.4 placeholder")
        reader = self._make_mock_reader(["text"])
        with patch("src.infrastructure.loaders.pdf_loader.PdfReader", return_value=reader):
            doc = PdfLoader().load(path)
        assert Path(doc.source).is_absolute()

    def test_none_page_text_handled(self, tmp_path: Path):
        path = tmp_path / "test.pdf"
        path.write_bytes(b"%PDF-1.4 placeholder")
        page = MagicMock()
        page.extract_text.return_value = None
        reader = MagicMock()
        reader.pages = [page]
        with patch("src.infrastructure.loaders.pdf_loader.PdfReader", return_value=reader):
            doc = PdfLoader().load(path)
        assert isinstance(doc.content, str)

    def test_wraps_exception_as_document_load_error(self, tmp_path: Path):
        path = tmp_path / "bad.pdf"
        path.write_bytes(b"not a pdf")
        with (
            patch(
                "src.infrastructure.loaders.pdf_loader.PdfReader",
                side_effect=ValueError("bad pdf"),
            ),
            pytest.raises(DocumentLoadError) as exc_info,
        ):
            PdfLoader().load(path)
        assert exc_info.value.cause is not None


# ── DocxLoader ─────────────────────────────────────────────────────────────────


class TestDocxLoader:
    def test_returns_document(self, docx_file: Path):
        assert isinstance(DocxLoader().load(docx_file), Document)

    def test_content_contains_paragraphs(self, docx_file: Path):
        doc = DocxLoader().load(docx_file)
        assert "first paragraph" in doc.content.lower()
        assert "second paragraph" in doc.content.lower()

    def test_metadata_filename(self, docx_file: Path):
        doc = DocxLoader().load(docx_file)
        assert doc.metadata["filename"] == "sample.docx"
        assert doc.metadata["extension"] == ".docx"
        assert doc.metadata["loader"] == "docx"

    def test_metadata_sections(self, docx_file: Path):
        doc = DocxLoader().load(docx_file)
        assert "Introduction" in doc.metadata["sections"]
        assert "Details" in doc.metadata["sections"]
        assert doc.metadata["section"] == "Introduction"

    def test_source_is_absolute(self, docx_file: Path):
        doc = DocxLoader().load(docx_file)
        assert Path(doc.source).is_absolute()

    def test_missing_file_raises_document_load_error(self, tmp_path: Path):
        with pytest.raises(DocumentLoadError):
            DocxLoader().load(tmp_path / "ghost.docx")


# ── PptxLoader ─────────────────────────────────────────────────────────────────


class TestPptxLoader:
    def test_returns_document(self, pptx_file: Path):
        assert isinstance(PptxLoader().load(pptx_file), Document)

    def test_content_contains_slide_text(self, pptx_file: Path):
        doc = PptxLoader().load(pptx_file)
        assert "Welcome to the deck." in doc.content
        assert "Second slide with more detail." in doc.content
        assert "---" in doc.content

    def test_metadata_filename(self, pptx_file: Path):
        doc = PptxLoader().load(pptx_file)
        assert doc.metadata["filename"] == "sample.pptx"
        assert doc.metadata["extension"] == ".pptx"
        assert doc.metadata["loader"] == "pptx"

    def test_metadata_sections(self, pptx_file: Path):
        doc = PptxLoader().load(pptx_file)
        assert doc.metadata["slide_count"] == 2
        assert "Introduction" in doc.metadata["sections"]
        assert "Details" in doc.metadata["sections"]
        assert doc.metadata["section"] == "Introduction"
        assert len(doc.metadata["slides"]) == 2
        assert all(isinstance(record, SlideRecord) for record in doc.metadata["slides"])
        assert doc.metadata["slides"][0].title == "Introduction"
        assert "Welcome to the deck." in doc.metadata["slides"][0].text
        assert doc.metadata["slides"][1].title == "Details"

    def test_blank_slide_without_title(self, blank_pptx_file: Path):
        doc = PptxLoader().load(blank_pptx_file)
        assert doc.metadata["slide_count"] == 1
        assert doc.metadata["sections"] == []
        assert doc.metadata["slides"] == []
        assert "section" not in doc.metadata

    def test_untitled_middle_slide_omitted_from_sections(self, tmp_path: Path):
        path = tmp_path / "mixed.pptx"
        prs = python_pptx.Presentation()
        first = prs.slides.add_slide(prs.slide_layouts[1])
        _set_slide_title_and_body(first, "Introduction", "Welcome.")

        untitled = prs.slides.add_slide(prs.slide_layouts[6])
        box = untitled.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        box.text_frame.text = "Agenda bullet without a slide title"

        last = prs.slides.add_slide(prs.slide_layouts[1])
        _set_slide_title_and_body(last, "Details", "More detail.")
        prs.save(str(path))

        doc = PptxLoader().load(path)
        assert doc.metadata["slide_count"] == 3
        assert doc.metadata["sections"] == ["Introduction", "Details"]
        assert doc.content.count("\n\n---\n\n") == 2
        assert "Agenda bullet without a slide title" in doc.content
        assert [slide.title for slide in doc.metadata["slides"]] == [
            "Introduction",
            None,
            "Details",
        ]

        from src.rag.chunking.section_chunker import SectionChunker

        chunks = SectionChunker().chunk(doc)
        by_section = {c.metadata.get(CHUNK_SECTION_KEY): c.text for c in chunks}
        assert "Welcome." in by_section["Introduction"]
        assert (
            "Agenda bullet without a slide title"
            in by_section["Agenda bullet without a slide title"]
        )
        assert "More detail." in by_section["Details"]
        assert all("slides" not in c.metadata for c in chunks)

    def test_agenda_slide_listing_titles_does_not_steal_pending(self, tmp_path: Path):
        path = tmp_path / "agenda.pptx"
        prs = python_pptx.Presentation()
        first = prs.slides.add_slide(prs.slide_layouts[1])
        _set_slide_title_and_body(first, "Introduction", "Welcome.")

        agenda = prs.slides.add_slide(prs.slide_layouts[6])
        box = agenda.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(3))
        box.text_frame.text = "Agenda\nDetails\nNext steps"

        last = prs.slides.add_slide(prs.slide_layouts[1])
        _set_slide_title_and_body(last, "Details", "More detail.")
        prs.save(str(path))

        from src.rag.chunking.section_chunker import SectionChunker

        doc = PptxLoader().load(path)
        chunks = SectionChunker().chunk(doc)
        by_section = {c.metadata.get(CHUNK_SECTION_KEY): c.text for c in chunks}
        assert "Welcome." in by_section["Introduction"]
        assert "Details" in by_section["Agenda"]
        assert "More detail." in by_section["Details"]
        assert by_section["Agenda"] != by_section["Details"]

    def test_intra_slide_hr_does_not_split_section_chunks(self, tmp_path: Path):
        path = tmp_path / "hr.pptx"
        prs = python_pptx.Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        _set_slide_title_and_body(slide, "Rules", "Before")
        # Extra shapes become "\n\n"-joined body parts, which can embed the
        # content separator sequence inside a single slide.
        hr = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(0.5))
        hr.text_frame.text = "---"
        after = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.5))
        after.text_frame.text = "After"
        prs.save(str(path))

        doc = PptxLoader().load(path)
        assert len(doc.metadata["slides"]) == 1
        assert "\n\n---\n\n" in doc.metadata["slides"][0].text

        from src.rag.chunking.section_chunker import SectionChunker

        chunks = SectionChunker().chunk(doc)
        assert len(chunks) == 1
        assert chunks[0].metadata[CHUNK_SECTION_KEY] == "Rules"
        assert "Before" in chunks[0].text and "After" in chunks[0].text

    def test_heading_like_body_does_not_steal_slide_sections(self, tmp_path: Path):
        path = tmp_path / "atx.pptx"
        prs = python_pptx.Presentation()
        first = prs.slides.add_slide(prs.slide_layouts[1])
        _set_slide_title_and_body(first, "Intro", "Welcome")
        key_points = first.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
        key_points.text_frame.text = "# Key Points"
        last = prs.slides.add_slide(prs.slide_layouts[1])
        _set_slide_title_and_body(last, "Details", "More detail.")
        prs.save(str(path))

        from src.rag.chunking.section_chunker import SectionChunker

        doc = PptxLoader().load(path)
        assert "# Key Points" in doc.content
        chunks = SectionChunker().chunk(doc)
        sections = {c.metadata.get(CHUNK_SECTION_KEY) for c in chunks}
        assert sections == {"Intro", "Details"}
        assert "Key Points" not in sections

    def test_source_is_absolute(self, pptx_file: Path):
        doc = PptxLoader().load(pptx_file)
        assert Path(doc.source).is_absolute()

    def test_missing_file_raises_document_load_error(self, tmp_path: Path):
        with pytest.raises(DocumentLoadError):
            PptxLoader().load(tmp_path / "ghost.pptx")


class TestPptxLoaderHelpers:
    def test_shape_text_without_text_frame(self):
        shape = MagicMock()
        shape.has_text_frame = False
        assert shape_text(shape) == ""

    def test_slide_title_missing_shape(self):
        slide = MagicMock()
        slide.shapes.title = None
        assert slide_title(slide) is None

    def test_slide_title_empty_text(self):
        slide = MagicMock()
        title_shape = MagicMock()
        title_shape.has_text_frame = True
        title_shape.text_frame.paragraphs = [MagicMock(text="   ")]
        slide.shapes.title = title_shape
        assert slide_title(slide) is None

    def test_slide_text_skips_empty_shapes(self):
        slide = MagicMock()
        text_shape = MagicMock()
        text_shape.has_text_frame = True
        text_shape.text_frame.paragraphs = [MagicMock(text="Slide body")]
        empty_shape = MagicMock()
        empty_shape.has_text_frame = False
        slide.shapes = [empty_shape, text_shape]
        assert slide_text(slide) == "Slide body"


# ── HtmlLoader ─────────────────────────────────────────────────────────────────


class TestHtmlLoader:
    def test_returns_document(self, html_file: Path):
        assert isinstance(HtmlLoader().load(html_file), Document)

    def test_main_content_preserved(self, html_file: Path):
        doc = HtmlLoader().load(html_file)
        assert "Main Content" in doc.content
        assert "main paragraph" in doc.content.lower()

    def test_boilerplate_stripped(self, html_file: Path):
        doc = HtmlLoader().load(html_file)
        assert "Navigation links" not in doc.content
        assert "Site Header" not in doc.content
        assert "Footer text" not in doc.content
        assert "alert" not in doc.content
        assert "color:red" not in doc.content

    def test_title_in_metadata(self, html_file: Path):
        doc = HtmlLoader().load(html_file)
        assert doc.metadata["title"] == "Test Page"

    def test_metadata_keys(self, html_file: Path):
        doc = HtmlLoader().load(html_file)
        assert doc.metadata["loader"] == "html"
        assert doc.metadata["extension"] == ".html"

    def test_latin1_fallback(self, tmp_path: Path):
        path = tmp_path / "latin1.html"
        path.write_bytes(b"<html><body><p>Caf\xe9</p></body></html>")
        doc = HtmlLoader().load(path)
        assert isinstance(doc.content, str)
        assert len(doc.content) > 0

    def test_missing_file_raises_document_load_error(self, tmp_path: Path):
        with pytest.raises(DocumentLoadError):
            HtmlLoader().load(tmp_path / "ghost.html")


# ── MarkdownLoader ─────────────────────────────────────────────────────────────


class TestMarkdownLoader:
    def test_returns_document(self, markdown_file: Path):
        assert isinstance(MarkdownLoader().load(markdown_file), Document)

    def test_raw_content_preserved(self, markdown_file: Path):
        doc = MarkdownLoader().load(markdown_file)
        assert "# Title" in doc.content
        assert "## Section One" in doc.content

    def test_headings_in_metadata(self, markdown_file: Path):
        doc = MarkdownLoader().load(markdown_file)
        assert doc.metadata["heading_count"] == 3
        assert "Title" in doc.metadata["headings"]
        assert "Section One" in doc.metadata["headings"]
        assert "Section Two" in doc.metadata["headings"]
        assert doc.metadata["section"] == "Title"

    def test_headings_skip_fenced_code_comments(self, tmp_path: Path):
        path = tmp_path / "fenced.md"
        path.write_text(
            "# Real\n\n```python\n# comment\nprint(1)\n```\n\n## Also\n",
            encoding="utf-8",
        )
        doc = MarkdownLoader().load(path)
        assert doc.metadata["headings"] == ["Real", "Also"]
        assert doc.metadata["heading_count"] == 2
        assert doc.metadata["section"] == "Real"

    def test_metadata_keys(self, markdown_file: Path):
        doc = MarkdownLoader().load(markdown_file)
        assert doc.metadata["loader"] == "markdown"
        assert doc.metadata["extension"] == ".md"

    def test_latin1_fallback(self, latin1_file: Path):
        doc = MarkdownLoader().load(latin1_file)
        assert isinstance(doc.content, str)
        assert len(doc.content) > 0

    def test_missing_file_raises_document_load_error(self, tmp_path: Path):
        with pytest.raises(DocumentLoadError):
            MarkdownLoader().load(tmp_path / "ghost.md")


# ── load_document factory ──────────────────────────────────────────────────────


class TestLoadDocument:
    def test_dispatches_to_markdown_loader(self, markdown_file: Path):
        doc = load_document(markdown_file)
        assert doc.metadata["loader"] == "markdown"

    def test_dispatches_to_html_loader(self, html_file: Path):
        doc = load_document(html_file)
        assert doc.metadata["loader"] == "html"

    def test_dispatches_to_docx_loader(self, docx_file: Path):
        doc = load_document(docx_file)
        assert doc.metadata["loader"] == "docx"

    def test_dispatches_to_pptx_loader(self, pptx_file: Path):
        doc = load_document(pptx_file)
        assert doc.metadata["loader"] == "pptx"

    def test_unsupported_extension_raises(self, tmp_path: Path):
        path = tmp_path / "file.xyz"
        path.write_text("content")
        with pytest.raises(DocumentLoadError, match="Unsupported"):
            load_document(path)

    def test_htm_extension_uses_html_loader(self, tmp_path: Path):
        path = tmp_path / "page.htm"
        path.write_text("<html><body><p>hello</p></body></html>")
        doc = load_document(path)
        assert doc.metadata["loader"] == "html"

    def test_dot_markdown_extension(self, tmp_path: Path):
        path = tmp_path / "readme.markdown"
        path.write_text("# Hello\n\nWorld.")
        doc = load_document(path)
        assert doc.metadata["loader"] == "markdown"

    def test_reload_settings_refreshes_layout_parser_routing(self, tmp_path: Path):
        """Env reloads must change whether PDFs route through the layout parser."""
        import src.infrastructure.loaders as loaders_mod
        from src.evals.e2e.technique_benchmark import temporary_config

        path = tmp_path / "report.pdf"
        path.write_bytes(b"%PDF-1.4 placeholder")

        reader = MagicMock()
        page = MagicMock()
        page.extract_text.return_value = "plain pdf text"
        reader.pages = [page]

        key = "PARSING__LAYOUT_PARSER__ENABLED"
        with temporary_config({key: "false"}):
            assert loaders_mod._settings().parsing.layout_parser.enabled is False
            with patch(
                "src.infrastructure.loaders.pdf_loader.PdfReader",
                return_value=reader,
            ):
                doc = load_document(path)
            assert doc.metadata["loader"] == "pdf"

        parsed = ParsedDocument(
            source=str(path.resolve()),
            content="layout text",
            metadata={"loader": "docling", "filename": "report.pdf"},
        )
        mock_parser = MagicMock()
        mock_parser.parse.return_value = parsed

        with temporary_config({key: "true"}):
            assert loaders_mod._settings().parsing.layout_parser.enabled is True
            with patch(
                "src.infrastructure.parsers.get_layout_parser",
                return_value=mock_parser,
            ):
                doc = load_document(path)
            mock_parser.parse.assert_called_once_with(path)
            assert doc.metadata["loader"] == "docling"
