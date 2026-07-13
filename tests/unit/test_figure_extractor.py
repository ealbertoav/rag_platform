"""T-230 — figure asset extraction and Chunk builders."""

from __future__ import annotations

import io
import logging
from collections.abc import Callable
from pathlib import Path
from types import SimpleNamespace
from typing import Any, cast
from unittest.mock import MagicMock, patch

import pptx as python_pptx
import pytest
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from src.core.constants import (
    ASSET_PATH_KEY,
    BBOX_KEY,
    CHUNK_PAGE_KEY,
    CHUNK_TYPE_FIGURE,
    CHUNK_TYPE_KEY,
    FIGURE_ID_KEY,
    MODALITY_FIGURE,
)
from src.core.exceptions import ConfigurationError, DocumentLoadError
from src.core.settings import FigureAssetSettings, ParsingSettings, Settings
from src.domain.entities.chunk import Chunk
from src.domain.entities.document import Document
from src.rag.ingestion.figure_extractor import (
    apply_figure_assets,
    build_figure_chunks,
    figure_chunk_id,
    is_figure_chunk,
)
from src.rag.ingestion.local_asset_store import LocalAssetStore

_EXTRACTOR = "src.rag.ingestion.figure_extractor"


def _settings(*, enabled: bool = True, store_dir: str | None = None) -> Settings:
    figure_assets = FigureAssetSettings(enabled=enabled)
    if store_dir is not None:
        figure_assets = figure_assets.model_copy(update={"store_dir": store_dir})
    return Settings(parsing=ParsingSettings(figure_assets=figure_assets))


def _png_bytes(color: tuple[int, int, int] = (255, 0, 0)) -> bytes:
    from PIL import Image

    image = Image.new("RGB", (8, 8), color=color)
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


def _save_png(fp: Any, *_args: Any, **_kwargs: Any) -> None:
    fp.write(_png_bytes())


def _save_empty(*_args: Any, **_kwargs: Any) -> None:
    return None


def _fake_pil_image(*, empty: bool = False) -> SimpleNamespace:
    return SimpleNamespace(save=_save_empty if empty else _save_png)


def _pptx_with_picture(tmp_path: Path) -> Path:
    png_path = tmp_path / "tiny.png"
    png_path.write_bytes(_png_bytes())
    pptx_path = tmp_path / "deck.pptx"
    presentation = python_pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(str(png_path), Inches(0.5), Inches(0.5), width=Inches(1))
    presentation.save(str(pptx_path))
    return pptx_path


def _pptx_with_grouped_picture(tmp_path: Path) -> Path:
    png_path = tmp_path / "grouped.png"
    png_path.write_bytes(_png_bytes((0, 0, 255)))
    pptx_path = tmp_path / "grouped.pptx"
    presentation = python_pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    group.shapes.add_picture(str(png_path), Inches(0.2), Inches(0.2), width=Inches(0.8))
    presentation.save(str(pptx_path))
    return pptx_path


def _pptx_with_nested_group_picture(tmp_path: Path) -> Path:
    png_path = tmp_path / "nested.png"
    png_path.write_bytes(_png_bytes((1, 2, 3)))
    pptx_path = tmp_path / "nested.pptx"
    presentation = python_pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    outer = slide.shapes.add_group_shape()
    inner = outer.shapes.add_group_shape()
    inner.shapes.add_picture(str(png_path), Inches(0.1), Inches(0.1), width=Inches(0.5))
    presentation.save(str(pptx_path))
    return pptx_path


def _pptx_text_only(tmp_path: Path) -> Path:
    pptx_path = tmp_path / "text-only.pptx"
    presentation = python_pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1)).text = "Hello"
    presentation.save(str(pptx_path))
    return pptx_path


def _apply_enabled_pptx_figures(tmp_path: Path, pptx_path: Path) -> list[dict[str, Any]]:
    """Run apply_figure_assets on a PPTX with assets enabled; return figures[]."""
    result = apply_figure_assets(
        Document(source=str(pptx_path.resolve()), content="slide text"),
        pptx_path,
        app_settings=_settings(enabled=True),
        store=LocalAssetStore(tmp_path / "assets"),
    )
    figures = result.metadata.get("figures")
    assert isinstance(figures, list)
    return figures


def _figure_pdf_document(
    tmp_path: Path,
    *,
    figures: list[Any] | None = None,
    extra_metadata: dict[str, Any] | None = None,
) -> tuple[Path, Document]:
    path = tmp_path / "doc.pdf"
    path.write_bytes(b"%PDF")
    metadata: dict[str, Any] = {
        "figures": figures if figures is not None else [{FIGURE_ID_KEY: "figure-1"}],
    }
    if extra_metadata:
        metadata.update(extra_metadata)
    return path, Document(source=str(path.resolve()), content="text", metadata=metadata)


def _apply_with_docling_pictures(
    document: Document,
    path: Path,
    store_dir: Path,
    pictures: list[Any],
    *,
    status_name: str | None = "SUCCESS",
    convert_error: BaseException | None = None,
) -> Document:
    converter = MagicMock()
    if convert_error is not None:
        converter.convert.side_effect = convert_error
    else:
        status = None if status_name is None else SimpleNamespace(name=status_name)
        converter.convert.return_value = SimpleNamespace(
            status=status,
            document=SimpleNamespace(pictures=list(pictures)),
        )
    with patch(f"{_EXTRACTOR}._create_picture_converter", return_value=converter):
        return apply_figure_assets(
            document,
            path,
            app_settings=_settings(enabled=True),
            store=LocalAssetStore(store_dir),
        )


class TestFigureChunkHelpers:
    def test_figure_chunk_id_stable(self) -> None:
        assert figure_chunk_id("/a.pdf", "figure-1") == figure_chunk_id("/a.pdf", "figure-1")
        assert figure_chunk_id("/a.pdf", "figure-1") != figure_chunk_id("/a.pdf", "figure-2")

    def test_is_figure_chunk(self) -> None:
        chunk = Chunk(
            document_id="d1",
            text="caption",
            metadata={CHUNK_TYPE_KEY: CHUNK_TYPE_FIGURE},
            modality=MODALITY_FIGURE,
            asset_path="/assets/f.png",
        )
        assert is_figure_chunk(chunk) is True
        assert is_figure_chunk(Chunk(document_id="d1", text="x")) is False


class TestBuildFigureChunks:
    def test_empty_without_figures(self) -> None:
        doc = Document(source="/doc.pdf", content="text")
        assert build_figure_chunks(doc) == []

    def test_skips_invalid_entries(self) -> None:
        doc = Document(
            source="/doc.pdf",
            content="text",
            metadata={
                "figures": [
                    "bad",
                    {FIGURE_ID_KEY: "figure-1"},  # no asset_path
                    {
                        FIGURE_ID_KEY: "figure-2",
                        ASSET_PATH_KEY: "/a/f2.png",
                        "caption": "  Chart  ",
                        CHUNK_PAGE_KEY: 3,
                        BBOX_KEY: [1.0, 2.0, 3.0, 4.0],
                    },
                ]
            },
        )
        chunks = build_figure_chunks(doc)
        assert len(chunks) == 1
        chunk = chunks[0]
        assert chunk.id == figure_chunk_id("/doc.pdf", "figure-2")
        assert chunk.text == "Chart"
        assert chunk.modality == MODALITY_FIGURE
        assert chunk.asset_path == "/a/f2.png"
        assert chunk.metadata[FIGURE_ID_KEY] == "figure-2"
        assert chunk.metadata[CHUNK_TYPE_KEY] == CHUNK_TYPE_FIGURE
        assert chunk.metadata[CHUNK_PAGE_KEY] == 3
        assert chunk.metadata[BBOX_KEY] == [1.0, 2.0, 3.0, 4.0]

    def test_default_text_without_caption(self) -> None:
        doc = Document(
            source="/doc.pdf",
            content="text",
            metadata={
                "figures": [{FIGURE_ID_KEY: "figure-1", ASSET_PATH_KEY: "/a/f1.png"}],
            },
        )
        chunks = build_figure_chunks(doc)
        assert chunks[0].text == "[figure]"


class TestApplyFigureAssetsDisabled:
    def test_noop_when_disabled(self, tmp_path: Path) -> None:
        path = tmp_path / "deck.pptx"
        path.write_bytes(b"not-used")
        doc = Document(source=str(path), content="slides")
        result = apply_figure_assets(doc, path, app_settings=_settings(enabled=False))
        assert result is doc

    def test_uses_default_settings_when_none(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ) -> None:
        monkeypatch.setenv("PARSING__FIGURE_ASSETS__ENABLED", "false")
        path = tmp_path / "x.pdf"
        path.write_text("x")
        doc = Document(source=str(path), content="text")
        # Reload settings via explicit Settings() path is covered by app_settings=None
        # with the live singleton still disabled by YAML default.
        result = apply_figure_assets(doc, path, app_settings=None)
        assert result is doc


class TestApplyPptxFigures:
    def test_persists_picture_and_sets_asset_path(self, tmp_path: Path) -> None:
        pptx_path = _pptx_with_picture(tmp_path)
        store = LocalAssetStore(tmp_path / "assets")
        doc = Document(source=str(pptx_path.resolve()), content="slide text")
        result = apply_figure_assets(
            doc,
            pptx_path,
            app_settings=_settings(enabled=True, store_dir=str(tmp_path / "assets")),
            store=store,
        )
        figures = result.metadata["figures"]
        assert len(figures) == 1
        assert figures[0][FIGURE_ID_KEY] == "figure-1"
        assert figures[0][CHUNK_PAGE_KEY] == 1
        asset = Path(figures[0][ASSET_PATH_KEY])
        assert asset.is_file()
        assert asset.read_bytes()

        chunks = build_figure_chunks(result)
        assert len(chunks) == 1
        assert chunks[0].asset_path == str(asset)
        assert chunks[0].metadata[FIGURE_ID_KEY] == "figure-1"

    def test_merges_existing_figure_metadata(self, tmp_path: Path) -> None:
        pptx_path = _pptx_with_picture(tmp_path)
        store = LocalAssetStore(tmp_path / "assets")
        doc = Document(
            source=str(pptx_path.resolve()),
            content="slide text",
            metadata={
                "figures": [
                    {FIGURE_ID_KEY: "figure-1", "caption": "Existing caption"},
                ]
            },
        )
        result = apply_figure_assets(
            doc,
            pptx_path,
            app_settings=_settings(enabled=True),
            store=store,
        )
        assert result.metadata["figures"][0]["caption"] == "Existing caption"
        assert ASSET_PATH_KEY in result.metadata["figures"][0]

    def test_persists_grouped_picture(self, tmp_path: Path) -> None:
        figures = _apply_enabled_pptx_figures(tmp_path, _pptx_with_grouped_picture(tmp_path))
        assert len(figures) == 1
        assert Path(figures[0][ASSET_PATH_KEY]).read_bytes() == _png_bytes((0, 0, 255))

    def test_persists_nested_group_picture(self, tmp_path: Path) -> None:
        figures = _apply_enabled_pptx_figures(tmp_path, _pptx_with_nested_group_picture(tmp_path))
        assert len(figures) == 1
        assert Path(figures[0][ASSET_PATH_KEY]).read_bytes() == _png_bytes((1, 2, 3))

    def test_soft_fails_store_error(self, tmp_path: Path, caplog: pytest.LogCaptureFixture) -> None:
        pptx_path = _pptx_with_picture(tmp_path)
        store = MagicMock(spec=LocalAssetStore)
        store.save.side_effect = OSError("disk full")
        doc = Document(source=str(pptx_path.resolve()), content="slide text")
        with caplog.at_level(logging.WARNING):
            result = apply_figure_assets(
                doc,
                pptx_path,
                app_settings=_settings(enabled=True),
                store=store,
            )
        assert ASSET_PATH_KEY not in result.metadata["figures"][0]
        assert "Failed to store PPTX figure" in caplog.text

    def test_skips_unreadable_picture(
        self,
        tmp_path: Path,
        caplog: pytest.LogCaptureFixture,
    ) -> None:
        pptx_path = _pptx_with_picture(tmp_path)
        doc = Document(source=str(pptx_path.resolve()), content="slide text")
        store = LocalAssetStore(tmp_path / "assets")

        class ExplodingImage:
            def __getattr__(self, name: str) -> object:
                raise RuntimeError("boom")

        exploding_shape = SimpleNamespace(
            shape_type=MSO_SHAPE_TYPE.PICTURE,
            image=ExplodingImage(),
        )
        good_shape = SimpleNamespace(
            shape_type=MSO_SHAPE_TYPE.PICTURE,
            image=SimpleNamespace(blob=_png_bytes((0, 255, 0)), ext="png"),
        )
        empty_shape = SimpleNamespace(
            shape_type=MSO_SHAPE_TYPE.PICTURE,
            image=SimpleNamespace(blob=b"", ext="png"),
        )
        text_shape = SimpleNamespace(shape_type=1)
        slide = SimpleNamespace(shapes=[exploding_shape, empty_shape, text_shape, good_shape])
        presentation = SimpleNamespace(slides=[slide])

        with (
            patch("pptx.Presentation", return_value=presentation),
            caplog.at_level(logging.WARNING),
        ):
            result = apply_figure_assets(
                doc,
                pptx_path,
                app_settings=_settings(enabled=True),
                store=store,
            )
        assert len(result.metadata["figures"]) == 1
        assert ASSET_PATH_KEY in result.metadata["figures"][0]
        assert "Skipping unreadable PPTX picture" in caplog.text


class TestApplyDoclingFigures:
    def test_noop_without_layout_figures(self, tmp_path: Path) -> None:
        path, doc = _figure_pdf_document(tmp_path, figures=[], extra_metadata={"loader": "docling"})
        result = apply_figure_assets(
            doc,
            path,
            app_settings=_settings(enabled=True),
            store=LocalAssetStore(tmp_path / "assets"),
        )
        assert result is doc

    def test_persists_docling_picture(
        self, tmp_path: Path, caplog: pytest.LogCaptureFixture
    ) -> None:
        path, doc = _figure_pdf_document(
            tmp_path,
            figures=[
                {
                    FIGURE_ID_KEY: "figure-1",
                    "caption": "Plot",
                    CHUNK_PAGE_KEY: 2,
                    BBOX_KEY: [0.0, 0.0, 10.0, 10.0],
                },
                "skip-me",
                {FIGURE_ID_KEY: "figure-2"},
            ],
            extra_metadata={"loader": "docling"},
        )
        picture = MagicMock()
        picture.get_image.return_value = _fake_pil_image()
        with caplog.at_level(logging.WARNING):
            result = _apply_with_docling_pictures(doc, path, tmp_path / "assets", [picture])

        figures = [f for f in result.metadata["figures"] if isinstance(f, dict)]
        assert figures[0][ASSET_PATH_KEY]
        assert Path(figures[0][ASSET_PATH_KEY]).is_file()
        # Second valid figure has no matching picture → no asset_path
        assert ASSET_PATH_KEY not in figures[1]
        assert "Docling returned 1 picture(s)" in caplog.text
        assert "No Docling picture for figure-2" in caplog.text
        chunks = build_figure_chunks(result)
        assert len(chunks) == 1
        assert chunks[0].text == "Plot"

    def test_aligns_pictures_past_non_dict_figure_entries(self, tmp_path: Path) -> None:
        path, doc = _figure_pdf_document(
            tmp_path,
            figures=[
                {FIGURE_ID_KEY: "figure-1"},
                "skip-me",
                {FIGURE_ID_KEY: "figure-2"},
            ],
        )
        first = MagicMock()
        first.get_image.return_value = SimpleNamespace(
            save=lambda fp, *_a, **_k: fp.write(_png_bytes((255, 0, 0)))
        )
        second = MagicMock()
        second.get_image.return_value = SimpleNamespace(
            save=lambda fp, *_a, **_k: fp.write(_png_bytes((0, 255, 0)))
        )
        result = _apply_with_docling_pictures(doc, path, tmp_path / "assets", [first, second])
        figures = [f for f in result.metadata["figures"] if isinstance(f, dict)]
        assert ASSET_PATH_KEY in figures[0]
        assert ASSET_PATH_KEY in figures[1]
        assert (
            Path(figures[0][ASSET_PATH_KEY]).read_bytes()
            != Path(figures[1][ASSET_PATH_KEY]).read_bytes()
        )

    def test_matches_pictures_by_page_and_bbox_when_order_differs(self, tmp_path: Path) -> None:
        path, doc = _figure_pdf_document(
            tmp_path,
            figures=[
                {
                    FIGURE_ID_KEY: "figure-1",
                    CHUNK_PAGE_KEY: 1,
                    BBOX_KEY: [0.0, 0.0, 10.0, 10.0],
                },
                {
                    FIGURE_ID_KEY: "figure-2",
                    CHUNK_PAGE_KEY: 2,
                    BBOX_KEY: [1.0, 2.0, 3.0, 4.0],
                },
            ],
        )
        first = MagicMock()
        first.prov = [SimpleNamespace(page_no=2, bbox=SimpleNamespace(l=1, t=2, r=3, b=4))]
        first.get_image.return_value = SimpleNamespace(
            save=lambda fp, *_a, **_k: fp.write(_png_bytes((0, 255, 0)))
        )
        second = MagicMock()
        second.prov = [SimpleNamespace(page_no=1, bbox=SimpleNamespace(l=0, t=0, r=10, b=10))]
        second.get_image.return_value = SimpleNamespace(
            save=lambda fp, *_a, **_k: fp.write(_png_bytes((255, 0, 0)))
        )
        # Intentionally reversed Docling picture order vs. layout figures[].
        result = _apply_with_docling_pictures(doc, path, tmp_path / "assets", [first, second])
        figures = result.metadata["figures"]
        assert Path(figures[0][ASSET_PATH_KEY]).read_bytes() == _png_bytes((255, 0, 0))
        assert Path(figures[1][ASSET_PATH_KEY]).read_bytes() == _png_bytes((0, 255, 0))

    def test_handles_missing_image_bytes(
        self, tmp_path: Path, caplog: pytest.LogCaptureFixture
    ) -> None:
        path, doc = _figure_pdf_document(tmp_path)
        picture = MagicMock()
        picture.get_image.return_value = None
        with caplog.at_level(logging.WARNING):
            result = _apply_with_docling_pictures(
                doc, path, tmp_path / "assets", [picture], status_name=None
            )
        assert ASSET_PATH_KEY not in result.metadata["figures"][0]
        assert "No image bytes" in caplog.text

    def test_handles_export_exception(
        self, tmp_path: Path, caplog: pytest.LogCaptureFixture
    ) -> None:
        path, doc = _figure_pdf_document(tmp_path)
        picture = MagicMock()
        picture.get_image.side_effect = RuntimeError("export failed")
        with caplog.at_level(logging.WARNING):
            result = _apply_with_docling_pictures(doc, path, tmp_path / "assets", [picture])
        assert ASSET_PATH_KEY not in result.metadata["figures"][0]
        assert "Failed to export Docling figure" in caplog.text

    def test_conversion_failure_status(
        self, tmp_path: Path, caplog: pytest.LogCaptureFixture
    ) -> None:
        path, doc = _figure_pdf_document(tmp_path)
        with caplog.at_level(logging.WARNING):
            result = _apply_with_docling_pictures(
                doc, path, tmp_path / "assets", [], status_name="FAILURE"
            )
        assert result is doc or result.metadata.get("figures") == doc.metadata["figures"]
        assert "Docling figure conversion failed" in caplog.text
        assert "Docling returned no pictures" in caplog.text

    def test_convert_raises_document_load_error(
        self,
        tmp_path: Path,
        caplog: pytest.LogCaptureFixture,
    ) -> None:
        path, doc = _figure_pdf_document(tmp_path)
        with caplog.at_level(logging.WARNING):
            result = _apply_with_docling_pictures(
                doc,
                path,
                tmp_path / "assets",
                [],
                convert_error=RuntimeError("boom"),
            )
        assert result is doc
        assert "Docling figure conversion failed" in caplog.text
        assert "Docling returned no pictures" in caplog.text

    def test_configuration_error_soft_fails(
        self,
        tmp_path: Path,
        caplog: pytest.LogCaptureFixture,
    ) -> None:
        path, doc = _figure_pdf_document(tmp_path)
        with (
            patch(
                f"{_EXTRACTOR}._create_picture_converter",
                side_effect=ConfigurationError("missing docling"),
            ),
            caplog.at_level(logging.WARNING),
        ):
            result = apply_figure_assets(
                doc,
                path,
                app_settings=_settings(enabled=True),
                store=LocalAssetStore(tmp_path / "assets"),
            )
        assert result is doc
        assert "misconfigured" in caplog.text

    def test_empty_pictures_list(self, tmp_path: Path, caplog: pytest.LogCaptureFixture) -> None:
        path, doc = _figure_pdf_document(
            tmp_path,
            figures=[{FIGURE_ID_KEY: "figure-1"}, {FIGURE_ID_KEY: "figure-2"}],
        )
        with caplog.at_level(logging.WARNING):
            result = _apply_with_docling_pictures(doc, path, tmp_path / "assets", [])
        assert result is doc
        assert ASSET_PATH_KEY not in doc.metadata["figures"][0]
        assert "Layout figures[] has 2 entries" in caplog.text
        assert "Docling returned no pictures" in caplog.text
        assert "continuing without asset_path" in caplog.text

    def test_empty_pictures_list_singular_wording(
        self, tmp_path: Path, caplog: pytest.LogCaptureFixture
    ) -> None:
        path, doc = _figure_pdf_document(tmp_path)
        with caplog.at_level(logging.WARNING):
            result = _apply_with_docling_pictures(doc, path, tmp_path / "assets", [])
        assert result is doc
        assert "Layout figures[] has 1 entry" in caplog.text
        assert "Docling returned no pictures" in caplog.text

    def test_none_pictures_attribute_treated_as_empty(
        self, tmp_path: Path, caplog: pytest.LogCaptureFixture
    ) -> None:
        path, doc = _figure_pdf_document(tmp_path)
        converter = MagicMock()
        converter.convert.return_value = SimpleNamespace(
            status=SimpleNamespace(name="SUCCESS"),
            document=SimpleNamespace(pictures=None),
        )
        with (
            patch(f"{_EXTRACTOR}._create_picture_converter", return_value=converter),
            caplog.at_level(logging.WARNING),
        ):
            result = apply_figure_assets(
                doc,
                path,
                app_settings=_settings(enabled=True),
                store=LocalAssetStore(tmp_path / "assets"),
            )
        assert result is doc
        assert "Docling returned no pictures" in caplog.text


class TestApplyFigureAssetsEdgeCases:
    def test_unsupported_extension(self, tmp_path: Path, caplog: pytest.LogCaptureFixture) -> None:
        path = tmp_path / "notes.md"
        path.write_text("# hi")
        doc = Document(source=str(path), content="hi")
        with caplog.at_level(logging.DEBUG):
            result = apply_figure_assets(
                doc,
                path,
                app_settings=_settings(enabled=True),
                store=LocalAssetStore(tmp_path / "assets"),
            )
        assert result is doc

    def test_unexpected_exception_soft_fails(
        self,
        tmp_path: Path,
        caplog: pytest.LogCaptureFixture,
    ) -> None:
        path, doc = _figure_pdf_document(tmp_path)
        with (
            patch(
                f"{_EXTRACTOR}._apply_docling_figures",
                side_effect=RuntimeError("surprise"),
            ),
            caplog.at_level(logging.WARNING),
        ):
            result = apply_figure_assets(
                doc,
                path,
                app_settings=_settings(enabled=True),
                store=LocalAssetStore(tmp_path / "assets"),
            )
        assert result is doc
        assert "Unexpected figure asset error" in caplog.text

    def test_pptx_open_failure(self, tmp_path: Path, caplog: pytest.LogCaptureFixture) -> None:
        path = tmp_path / "broken.pptx"
        path.write_bytes(b"not-a-pptx")
        doc = Document(source=str(path), content="")
        with caplog.at_level(logging.WARNING):
            result = apply_figure_assets(
                doc,
                path,
                app_settings=_settings(enabled=True),
                store=LocalAssetStore(tmp_path / "assets"),
            )
        assert result is doc
        assert "Figure asset extraction failed" in caplog.text

    def test_pptx_import_error(self, tmp_path: Path, caplog: pytest.LogCaptureFixture) -> None:
        path = tmp_path / "deck.pptx"
        path.write_bytes(b"x")
        doc = Document(source=str(path), content="")
        with (
            patch(
                f"{_EXTRACTOR}._extract_pptx_picture_bytes",
                side_effect=ConfigurationError("PPTX figure extraction requires python-pptx"),
            ),
            caplog.at_level(logging.WARNING),
        ):
            result = apply_figure_assets(
                doc,
                path,
                app_settings=_settings(enabled=True),
                store=LocalAssetStore(tmp_path / "assets"),
            )
        assert result is doc
        assert "misconfigured" in caplog.text

    def test_create_picture_converter_delegates_to_layout_converter(self) -> None:
        from src.rag.ingestion import figure_extractor as mod

        sentinel = object()
        with patch(
            "src.infrastructure.parsers.docling_parser.create_docling_converter",
            return_value=sentinel,
        ) as create:
            assert mod._create_picture_converter() is sentinel
        create.assert_called_once_with()

    def test_create_picture_converter_propagates_configuration_error(self) -> None:
        from src.rag.ingestion import figure_extractor as mod

        with (
            patch(
                "src.infrastructure.parsers.docling_parser.create_docling_converter",
                side_effect=ConfigurationError("missing docling"),
            ),
            pytest.raises(ConfigurationError, match="docling"),
        ):
            mod._create_picture_converter()

    def test_extract_pptx_import_error_directly(self, tmp_path: Path) -> None:
        from src.rag.ingestion import figure_extractor as mod

        path = tmp_path / "deck.pptx"
        path.write_bytes(b"x")
        real_import = __import__

        def fake_import(name: str, *args: Any, **kwargs: Any) -> Any:
            if name == "pptx" or name.startswith("pptx."):
                raise ImportError("missing pptx")
            return real_import(name, *args, **kwargs)

        with (
            patch("builtins.__import__", side_effect=fake_import),
            pytest.raises(ConfigurationError, match="python-pptx"),
        ):
            mod._extract_pptx_picture_bytes(path)

    def test_picture_to_png_empty_buffer(
        self, tmp_path: Path, caplog: pytest.LogCaptureFixture
    ) -> None:
        path, doc = _figure_pdf_document(tmp_path)
        picture = MagicMock()
        picture.get_image.return_value = _fake_pil_image(empty=True)
        with caplog.at_level(logging.WARNING):
            result = _apply_with_docling_pictures(doc, path, tmp_path / "assets", [picture])
        assert ASSET_PATH_KEY not in result.metadata["figures"][0]
        assert "No image bytes" in caplog.text

    def test_pptx_without_pictures(self, tmp_path: Path) -> None:
        pptx_path = _pptx_text_only(tmp_path)
        doc = Document(source=str(pptx_path.resolve()), content="Hello")
        result = apply_figure_assets(
            doc,
            pptx_path,
            app_settings=_settings(enabled=True),
            store=LocalAssetStore(tmp_path / "assets"),
        )
        assert result is doc

    def test_pptx_warns_when_figures_exist_but_no_pictures(
        self, tmp_path: Path, caplog: pytest.LogCaptureFixture
    ) -> None:
        pptx_path = _pptx_text_only(tmp_path)
        doc = Document(
            source=str(pptx_path.resolve()),
            content="Hello",
            metadata={
                "figures": [
                    {FIGURE_ID_KEY: "figure-1"},
                    {"caption": "no-id-ignored"},
                    {FIGURE_ID_KEY: "figure-2", ASSET_PATH_KEY: "/already/there.png"},
                    "bad-entry",
                ]
            },
        )
        with caplog.at_level(logging.WARNING):
            result = apply_figure_assets(
                doc,
                pptx_path,
                app_settings=_settings(enabled=True),
                store=LocalAssetStore(tmp_path / "assets"),
            )
        assert result is doc
        assert "1 figure(s) without asset_path" in caplog.text
        assert "no extractable pictures" in caplog.text
        assert "continuing without assets" in caplog.text

    def test_figure_entry_needs_asset_helpers(self) -> None:
        from src.rag.ingestion.figure_extractor import _figure_entry_needs_asset

        assert _figure_entry_needs_asset({}) is False
        assert _figure_entry_needs_asset({FIGURE_ID_KEY: "figure-1"}) is True
        assert (
            _figure_entry_needs_asset({FIGURE_ID_KEY: "figure-1", ASSET_PATH_KEY: "/a/f.png"})
            is False
        )
        assert (
            _figure_entry_needs_asset({FIGURE_ID_KEY: "figure-1", "asset_path": "/a/f.png"})
            is False
        )

    def test_docling_assigns_default_figure_id(self, tmp_path: Path) -> None:
        path, doc = _figure_pdf_document(tmp_path, figures=[{"caption": "No id"}])
        picture = MagicMock()
        picture.get_image.return_value = _fake_pil_image()
        result = _apply_with_docling_pictures(doc, path, tmp_path / "assets", [picture])
        assert result.metadata["figures"][0][FIGURE_ID_KEY] == "figure-1"
        assert ASSET_PATH_KEY in result.metadata["figures"][0]

    def test_ext_none_defaults_png(self, tmp_path: Path) -> None:
        from src.rag.ingestion import figure_extractor as mod

        shape = SimpleNamespace(
            shape_type=MSO_SHAPE_TYPE.PICTURE,
            image=SimpleNamespace(blob=_png_bytes(), ext=None),
        )
        presentation = SimpleNamespace(slides=[SimpleNamespace(shapes=[shape])])
        with patch("pptx.Presentation", return_value=presentation):
            pictures = mod._extract_pptx_picture_bytes(tmp_path / "x.pptx")
        assert len(pictures) == 1
        assert pictures[0][1] == "png"
        assert pictures[0][2] == 1

    def test_convert_raises_configuration_error(
        self,
        tmp_path: Path,
        caplog: pytest.LogCaptureFixture,
    ) -> None:
        path, doc = _figure_pdf_document(tmp_path)
        with caplog.at_level(logging.WARNING):
            result = _apply_with_docling_pictures(
                doc,
                path,
                tmp_path / "assets",
                [],
                convert_error=ConfigurationError("nested config"),
            )
        assert result is doc
        assert "misconfigured" in caplog.text


def _docx_with_picture(tmp_path: Path) -> Path:
    """Build a minimal DOCX containing one embedded PNG."""
    from docx import Document as DocxDocument
    from docx.shared import Inches as DocxInches

    png_path = tmp_path / "embed.png"
    png_path.write_bytes(_png_bytes((0, 128, 255)))
    docx_path = tmp_path / "report.docx"
    document = DocxDocument()
    document.add_paragraph("Figure below")
    document.add_picture(str(png_path), width=DocxInches(1))
    document.save(str(docx_path))
    return docx_path


def _figure_docx_document(
    tmp_path: Path,
    *,
    figures: list[Any] | None = None,
    with_picture: bool = True,
) -> tuple[Path, Document]:
    path = _docx_with_picture(tmp_path) if with_picture else (tmp_path / "plain.docx")
    if not with_picture:
        from docx import Document as DocxDocument

        DocxDocument().save(str(path))
    metadata: dict[str, Any] = {
        "figures": figures if figures is not None else [{FIGURE_ID_KEY: "figure-1"}],
        "loader": "docling",
    }
    return path, Document(source=str(path.resolve()), content="text", metadata=metadata)


class TestDocxFigureFallback:
    def test_docx_fallback_when_docling_has_no_image_bytes(self, tmp_path: Path) -> None:
        path, doc = _figure_docx_document(tmp_path)
        picture = MagicMock()
        picture.get_image.return_value = None
        result = _apply_with_docling_pictures(doc, path, tmp_path / "assets", [picture])
        figures = result.metadata["figures"]
        assert ASSET_PATH_KEY in figures[0]
        asset = Path(figures[0][ASSET_PATH_KEY])
        assert asset.is_file()
        assert asset.read_bytes()

    def test_docx_aligns_fallback_blobs_by_figure_slot(self, tmp_path: Path) -> None:
        """After a successful Docling export, DOCX fallback must use the same slot index."""
        path, doc = _figure_docx_document(
            tmp_path,
            figures=[{FIGURE_ID_KEY: "figure-1"}, {FIGURE_ID_KEY: "figure-2"}],
        )
        first = MagicMock()
        first.get_image.return_value = _fake_pil_image()
        second = MagicMock()
        second.get_image.return_value = None
        blob_a = _png_bytes((10, 20, 30))
        blob_b = _png_bytes((40, 50, 60))
        with patch(
            f"{_EXTRACTOR}._docx_fallback_blobs",
            return_value=[(blob_a, "png"), (blob_b, "png")],
        ):
            result = _apply_with_docling_pictures(doc, path, tmp_path / "assets", [first, second])
        figures = result.metadata["figures"]
        assert Path(figures[0][ASSET_PATH_KEY]).read_bytes() == _png_bytes()
        assert Path(figures[1][ASSET_PATH_KEY]).read_bytes() == blob_b

    def test_docx_fallback_when_docling_conversion_fails(
        self, tmp_path: Path, caplog: pytest.LogCaptureFixture
    ) -> None:
        path, doc = _figure_docx_document(tmp_path)
        with caplog.at_level(logging.WARNING):
            result = _apply_with_docling_pictures(
                doc,
                path,
                tmp_path / "assets",
                [],
                convert_error=RuntimeError("boom"),
            )
        assert ASSET_PATH_KEY in result.metadata["figures"][0]
        assert Path(result.metadata["figures"][0][ASSET_PATH_KEY]).is_file()
        assert "Docling figure conversion failed" in caplog.text
        assert "trying embedded-image fallback" in caplog.text

    def test_docx_fallback_when_docling_unavailable(
        self, tmp_path: Path, caplog: pytest.LogCaptureFixture
    ) -> None:
        path, doc = _figure_docx_document(tmp_path)
        with (
            patch(
                f"{_EXTRACTOR}._create_picture_converter",
                side_effect=ConfigurationError("missing docling"),
            ),
            caplog.at_level(logging.WARNING),
        ):
            result = apply_figure_assets(
                doc,
                path,
                app_settings=_settings(enabled=True),
                store=LocalAssetStore(tmp_path / "assets"),
            )
        assert ASSET_PATH_KEY in result.metadata["figures"][0]
        assert Path(result.metadata["figures"][0][ASSET_PATH_KEY]).is_file()
        assert "Docling unavailable" in caplog.text
        assert "python-docx embedded-image fallback" in caplog.text

    def test_docx_fallback_when_docling_returns_no_pictures(self, tmp_path: Path) -> None:
        path, doc = _figure_docx_document(tmp_path)
        result = _apply_with_docling_pictures(doc, path, tmp_path / "assets", [])
        assert ASSET_PATH_KEY in result.metadata["figures"][0]
        assert Path(result.metadata["figures"][0][ASSET_PATH_KEY]).is_file()

    def test_docx_fallback_mismatch_warning(
        self, tmp_path: Path, caplog: pytest.LogCaptureFixture
    ) -> None:
        path, doc = _figure_docx_document(
            tmp_path,
            figures=[{FIGURE_ID_KEY: "figure-1"}, {FIGURE_ID_KEY: "figure-2"}],
        )
        with caplog.at_level(logging.WARNING):
            result = _apply_with_docling_pictures(doc, path, tmp_path / "assets", [])
        assert ASSET_PATH_KEY in result.metadata["figures"][0]
        assert ASSET_PATH_KEY not in result.metadata["figures"][1]
        assert "DOCX embedded images returned 1 picture(s)" in caplog.text
        assert "No Docling picture for figure-2" in caplog.text

    def test_docx_open_failure_soft_falls_back_empty(
        self, tmp_path: Path, caplog: pytest.LogCaptureFixture
    ) -> None:
        from src.rag.ingestion import figure_extractor as mod

        path = tmp_path / "broken.docx"
        path.write_bytes(b"not-a-docx")
        doc = Document(
            source=str(path.resolve()),
            content="x",
            metadata={"figures": [{FIGURE_ID_KEY: "figure-1"}]},
        )
        with (
            patch(f"{_EXTRACTOR}._load_docling_pictures", return_value=[]),
            caplog.at_level(logging.WARNING),
        ):
            result = apply_figure_assets(
                doc,
                path,
                app_settings=_settings(enabled=True),
                store=LocalAssetStore(tmp_path / "assets"),
            )
        assert result is doc
        assert "DOCX embedded-image fallback unavailable" in caplog.text
        assert "Docling returned no pictures" in caplog.text
        with pytest.raises(DocumentLoadError, match="Cannot open DOCX"):
            mod._extract_docx_picture_bytes(path)

    def test_extract_docx_import_error(self, tmp_path: Path) -> None:
        from src.rag.ingestion import figure_extractor as mod

        path = tmp_path / "x.docx"
        path.write_bytes(b"x")
        real_import = __import__

        def fake_import(name: str, *args: Any, **kwargs: Any) -> Any:
            if name == "docx" or name.startswith("docx."):
                raise ImportError("missing docx")
            return real_import(name, *args, **kwargs)

        with (
            patch("builtins.__import__", side_effect=fake_import),
            pytest.raises(ConfigurationError, match="python-docx"),
        ):
            mod._extract_docx_picture_bytes(path)

    def test_extract_docx_skips_unreadable_and_empty_parts(
        self, tmp_path: Path, caplog: pytest.LogCaptureFixture
    ) -> None:
        from docx.oxml.ns import qn

        from src.rag.ingestion import figure_extractor as mod

        path = _docx_with_picture(tmp_path)
        image_reltype = "officeDocument/relationships/image"
        embed_attr = qn("r:embed")
        blip_tag = qn("a:blip")

        class FakeBlip:
            def __init__(self, rid: str | None) -> None:
                self._rid = rid

            def get(self, attr: str) -> str | None:
                if attr == embed_attr:
                    return self._rid
                return None

        class BoomRel:
            reltype = image_reltype

            def __getattr__(self, name: str) -> object:
                if name == "target_part":
                    raise RuntimeError("bad part")
                raise AttributeError(name)

        class EmptyRel:
            reltype = image_reltype
            target_part = SimpleNamespace(blob=b"", content_type="image/png")

        class GoodRel:
            reltype = image_reltype
            target_part = SimpleNamespace(blob=_png_bytes(), content_type="image/jpeg")

        class NonImageRel:
            reltype = "officeDocument/relationships/hyperlink"
            target_part = SimpleNamespace(blob=_png_bytes((9, 9, 9)), content_type="image/png")

        class FakeBody:
            @staticmethod
            def iter(tag: str) -> list[FakeBlip]:
                assert tag == blip_tag
                return [
                    FakeBlip("rIdBoom"),
                    FakeBlip("rIdEmpty"),
                    FakeBlip("rIdMissing"),
                    FakeBlip(None),
                    FakeBlip("rIdLink"),
                    FakeBlip("rIdGood"),
                    FakeBlip("rIdGood"),  # repeated embed still yields a slot
                ]

        fake_doc = SimpleNamespace(
            element=SimpleNamespace(body=FakeBody()),
            part=SimpleNamespace(
                rels={
                    "rIdBoom": BoomRel(),
                    "rIdEmpty": EmptyRel(),
                    "rIdLink": NonImageRel(),
                    "rIdGood": GoodRel(),
                }
            ),
        )
        with (
            patch("docx.Document", return_value=fake_doc),
            caplog.at_level(logging.WARNING),
        ):
            pictures = mod._extract_docx_picture_bytes(path)
        assert len(pictures) == 2
        assert pictures[0][1] == "jpg"
        assert pictures[1][1] == "jpg"
        assert "Skipping unreadable DOCX embedded image" in caplog.text
        assert "missing relationship" in caplog.text

    def test_extract_docx_returns_empty_without_body(self, tmp_path: Path) -> None:
        from src.rag.ingestion import figure_extractor as mod

        path = tmp_path / "nobody.docx"
        path.write_bytes(b"x")
        fake_doc = SimpleNamespace(
            element=SimpleNamespace(body=None),
            part=SimpleNamespace(rels={}),
        )
        with patch("docx.Document", return_value=fake_doc):
            assert mod._extract_docx_picture_bytes(path) == []

    def test_extract_docx_uses_body_order_not_rels_order(self, tmp_path: Path) -> None:
        from docx import Document as DocxDocument
        from docx.shared import Inches as DocxInches

        from src.rag.ingestion import figure_extractor as mod

        first = tmp_path / "first.png"
        second = tmp_path / "second.png"
        first.write_bytes(_png_bytes((10, 20, 30)))
        second.write_bytes(_png_bytes((40, 50, 60)))
        docx_path = tmp_path / "ordered.docx"
        document = DocxDocument()
        document.add_picture(str(first), width=DocxInches(0.5))
        document.add_picture(str(second), width=DocxInches(0.5))
        document.save(str(docx_path))

        pictures = mod._extract_docx_picture_bytes(docx_path)
        assert len(pictures) == 2
        assert pictures[0][0] == _png_bytes((10, 20, 30))
        assert pictures[1][0] == _png_bytes((40, 50, 60))

    def test_extension_from_image_content_type_helpers(self) -> None:
        from src.rag.ingestion.figure_extractor import _extension_from_image_content_type

        assert _extension_from_image_content_type("image/png") == "png"
        assert _extension_from_image_content_type("image/jpeg") == "jpg"
        assert _extension_from_image_content_type("image/x-emf") == "emf"
        assert _extension_from_image_content_type("image/x-icon") == "icon"
        assert _extension_from_image_content_type("image/svg+xml") == "svg+xml"
        assert _extension_from_image_content_type("") == "png"
        assert _extension_from_image_content_type("application/octet-stream") == "octet-stream"
        assert _extension_from_image_content_type("noslash") == "png"

    def test_docx_fallback_blobs_ignores_non_docx(self, tmp_path: Path) -> None:
        from src.rag.ingestion.figure_extractor import _docx_fallback_blobs

        assert _docx_fallback_blobs(tmp_path / "doc.pdf") == []

    def test_docx_fallback_unexpected_error_returns_empty(
        self, tmp_path: Path, caplog: pytest.LogCaptureFixture
    ) -> None:
        from src.rag.ingestion import figure_extractor as mod

        path = tmp_path / "report.docx"
        path.write_bytes(b"x")
        with (
            patch(
                f"{_EXTRACTOR}._extract_docx_picture_bytes",
                side_effect=RuntimeError("surprise"),
            ),
            caplog.at_level(logging.WARNING),
        ):
            assert mod._docx_fallback_blobs(path) == []
        assert "DOCX embedded-image fallback failed" in caplog.text

    def test_docx_prefers_docling_bytes_over_fallback(self, tmp_path: Path) -> None:
        path, doc = _figure_docx_document(tmp_path)
        picture = MagicMock()
        picture.get_image.return_value = _fake_pil_image()
        result = _apply_with_docling_pictures(doc, path, tmp_path / "assets", [picture])
        asset = Path(result.metadata["figures"][0][ASSET_PATH_KEY])
        # Fake PIL writes the standard red PNG from _png_bytes(), not the embedded blue one.
        assert asset.read_bytes() == _png_bytes()

    def test_docx_fallback_after_docling_export_exception(
        self, tmp_path: Path, caplog: pytest.LogCaptureFixture
    ) -> None:
        path, doc = _figure_docx_document(tmp_path)
        picture = MagicMock()
        picture.get_image.side_effect = RuntimeError("export failed")
        with caplog.at_level(logging.WARNING):
            result = _apply_with_docling_pictures(doc, path, tmp_path / "assets", [picture])
        assert ASSET_PATH_KEY in result.metadata["figures"][0]
        assert "Failed to export Docling figure" in caplog.text
        assert Path(result.metadata["figures"][0][ASSET_PATH_KEY]).is_file()

    def test_store_save_failure_after_docx_blob(
        self, tmp_path: Path, caplog: pytest.LogCaptureFixture
    ) -> None:
        path, doc = _figure_docx_document(tmp_path)
        store = MagicMock(spec=LocalAssetStore)
        store.save.side_effect = OSError("disk full")
        with (
            patch(f"{_EXTRACTOR}._load_docling_pictures", return_value=[]),
            caplog.at_level(logging.WARNING),
        ):
            result = apply_figure_assets(
                doc,
                path,
                app_settings=_settings(enabled=True),
                store=store,
            )
        assert ASSET_PATH_KEY not in result.metadata["figures"][0]
        assert "Failed to export Docling figure" in caplog.text

    def test_docx_fallback_configuration_error_propagates(
        self, tmp_path: Path, caplog: pytest.LogCaptureFixture
    ) -> None:
        path, doc = _figure_docx_document(tmp_path)
        with (
            patch(
                f"{_EXTRACTOR}._extract_docx_picture_bytes",
                side_effect=ConfigurationError("DOCX figure fallback requires python-docx"),
            ),
            patch(f"{_EXTRACTOR}._load_docling_pictures", return_value=[]),
            caplog.at_level(logging.WARNING),
        ):
            result = apply_figure_assets(
                doc,
                path,
                app_settings=_settings(enabled=True),
                store=LocalAssetStore(tmp_path / "assets"),
            )
        assert result is doc
        assert "misconfigured" in caplog.text


class TestPictureMatchingHelpers:
    def test_picture_page_and_bbox_edge_cases(self) -> None:
        from src.rag.ingestion.figure_extractor import _picture_bbox, _picture_page

        picture_page = cast(Callable[[object], int | None], _picture_page)
        picture_bbox = cast(Callable[[object], list[float] | None], _picture_bbox)

        assert picture_page(SimpleNamespace(prov=None)) is None
        assert picture_page(SimpleNamespace(prov=[])) is None
        assert picture_page(SimpleNamespace(prov=[SimpleNamespace()])) is None
        assert picture_page(SimpleNamespace(prov=[SimpleNamespace(page_no="bad")])) is None
        assert picture_page(SimpleNamespace(prov=[SimpleNamespace(page_no=3)])) == 3

        assert picture_bbox(SimpleNamespace(prov=None)) is None
        assert picture_bbox(SimpleNamespace(prov=[])) is None
        assert picture_bbox(SimpleNamespace(prov=[SimpleNamespace(bbox=None)])) is None
        assert picture_bbox(SimpleNamespace(prov=[object()])) is None
        assert (
            picture_bbox(SimpleNamespace(prov=[SimpleNamespace(bbox=SimpleNamespace(l="x"))]))
            is None
        )
        assert picture_bbox(
            SimpleNamespace(prov=[SimpleNamespace(bbox=SimpleNamespace(l=1, t=2, r=3, b=4))])
        ) == [1.0, 2.0, 3.0, 4.0]

    def test_bboxes_equal_and_match_picture_index(self) -> None:
        from src.rag.ingestion.figure_extractor import _bboxes_equal, _match_picture_index

        match_picture_index = cast(
            Callable[
                [dict[str, Any], list[tuple[object, object]], set[int], int],
                int | None,
            ],
            _match_picture_index,
        )

        assert _bboxes_equal([0, 0, 1, 1], [0, 0, 1, 1]) is True
        assert _bboxes_equal([0, 0, 1], [0, 0, 1, 1]) is False
        assert _bboxes_equal([0, 0, 1, 1], [0, 0, 1, 2]) is False

        first = SimpleNamespace(
            prov=[SimpleNamespace(page_no=1, bbox=SimpleNamespace(l=0, t=0, r=1, b=1))]
        )
        second = SimpleNamespace(
            prov=[SimpleNamespace(page_no=2, bbox=SimpleNamespace(l=5, t=5, r=6, b=6))]
        )
        pictures: list[tuple[object, object]] = [(first, object()), (second, object())]

        assert match_picture_index({}, [], set(), 0) is None
        assert (
            match_picture_index(
                {CHUNK_PAGE_KEY: 2, BBOX_KEY: [5, 5, 6, 6]},
                pictures,
                set(),
                0,
            )
            == 1
        )
        # Provenance hit already consumed → skip and fall back to unused slot.
        assert (
            match_picture_index(
                {CHUNK_PAGE_KEY: 2, BBOX_KEY: [5, 5, 6, 6]},
                pictures,
                {1},
                0,
            )
            == 0
        )
        # Page matches, but bbox does not → fall back.
        assert (
            match_picture_index(
                {CHUNK_PAGE_KEY: 2, BBOX_KEY: [9, 9, 9, 9]},
                pictures,
                set(),
                0,
            )
            == 0
        )
        # No page match → fall back to document order.
        assert (
            match_picture_index(
                {CHUNK_PAGE_KEY: 99, BBOX_KEY: [5, 5, 6, 6]},
                pictures,
                set(),
                0,
            )
            == 0
        )
        # Invalid provenance falls back to document-order slot.
        assert (
            match_picture_index(
                {CHUNK_PAGE_KEY: "bad", BBOX_KEY: [5, 5, 6, 6]},
                pictures,
                set(),
                0,
            )
            == 0
        )
        # When preferred fallback is already used, take the next unused index.
        assert match_picture_index({}, pictures, {0}, 0) == 1
        assert match_picture_index({}, pictures, {0, 1}, 0) is None

    def test_iter_pptx_shapes_skips_group_without_shapes(self) -> None:
        from src.rag.ingestion.figure_extractor import _iter_pptx_shapes

        group = SimpleNamespace(shape_type=MSO_SHAPE_TYPE.GROUP, shapes=None)
        picture = SimpleNamespace(shape_type=MSO_SHAPE_TYPE.PICTURE)
        assert list(_iter_pptx_shapes([group, picture])) == [picture]
